"""
Raphael Universal File Processor.

Supported types and actions (non-AI operations only — AI analysis is handled by the LLM):

  image   → resize, convert, compress, crop, info
  pdf     → extract_text, info, to_word, extract_pages
  docx    → extract_text, info
  text    → word_count, info
  csv     → filter, sort, convert, stats, info
  excel   → filter, convert, stats, info
  json    → validate, format, extract, convert
  code    → run, info
  audio   → info
  video   → trim, extract_audio, extract_frame, info, compress
  archive → list, extract
  pptx    → extract_text, info
"""

import logging
import re
import json
import subprocess
from pathlib import Path

logger = logging.getLogger(__name__)


# ── File type detection ──────────────────────────────────────────────

def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "svg", "ico"}
    video_exts = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}
    audio_exts = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "wma", "opus"}
    code_exts = {"py", "js", "ts", "jsx", "tsx", "html", "css", "java", "c",
                 "cpp", "cs", "go", "rs", "rb", "php", "swift", "kt", "sh",
                 "bash", "ps1", "lua", "r", "m", "sql", "yaml", "toml"}
    archive_exts = {"zip", "rar", "tar", "gz", "7z", "bz2", "xz"}

    if ext in image_exts:   return "image"
    if ext in video_exts:   return "video"
    if ext in audio_exts:   return "audio"
    if ext in code_exts:    return "code"
    if ext in archive_exts: return "archive"
    if ext == "pdf":        return "pdf"
    if ext in ("docx", "doc"):  return "docx"
    if ext in ("txt", "md", "rst", "log"): return "text"
    if ext in ("csv", "tsv"):   return "csv"
    if ext in ("xlsx", "xls", "ods"): return "excel"
    if ext == "json":       return "json"
    if ext == "xml":        return "xml"
    if ext in ("pptx", "ppt"):  return "pptx"
    return "unknown"


# ── Helpers ─────────────────────────────────────────────────────────

def _file_size_str(path: Path) -> str:
    size = path.stat().st_size
    if size < 1024:
        return f"{size} B"
    if size < 1024 ** 2:
        return f"{size / 1024:.1f} KB"
    if size < 1024 ** 3:
        return f"{size / 1024 ** 2:.1f} MB"
    return f"{size / 1024 ** 3:.1f} GB"


def _output_path(src: Path, suffix: str, new_ext: str | None = None) -> Path:
    ext = new_ext or src.suffix
    name = f"{src.stem}_{suffix}{ext}"
    return src.parent / name


# ── Image operations ────────────────────────────────────────────────

def _process_image(path: Path, action: str, params: dict) -> str:
    """Image operations (no AI — resize, convert, compress, crop, info)."""
    try:
        from PIL import Image
    except ImportError:
        return "Pillow is not installed. Run: pip install Pillow"

    action = action or "info"

    # AI-dependent actions — return extracted content for the LLM to process
    if action in ("describe", "ocr", "analyze", "read", "extract_text"):
        return (
            f"AI image analysis is handled by the LLM. "
            f"The image is at: {path} ({_file_size_str(path)}). "
            f"If you need text extracted, use the 'info' action to get details."
        )

    if action == "resize":
        width = int(params.get("width", 0))
        height = int(params.get("height", 0))
        scale = float(params.get("scale", 0))
        try:
            img = Image.open(path)
            w, h = img.size
            if scale:
                new_size = (int(w * scale), int(h * scale))
            elif width and height:
                new_size = (width, height)
            elif width:
                new_size = (width, int(h * width / w))
            elif height:
                new_size = (int(w * height / h), height)
            else:
                return "Specify width, height, or scale."
            out = _output_path(path, f"resized_{new_size[0]}x{new_size[1]}")
            img.resize(new_size, Image.LANCZOS).save(out)  # type: ignore[attr-defined]
            return f"Resized from {w}x{h} to {new_size[0]}x{new_size[1]}. Saved: {out.name}"
        except Exception as e:
            return f"Resize failed: {e}"

    if action == "convert":
        fmt = params.get("format", "png").lower().strip(".")
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG",
                   "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}
        pil_fmt = fmt_map.get(fmt, fmt.upper())
        try:
            img = Image.open(path).convert("RGB") if fmt == "jpg" else Image.open(path)  # type: ignore[assignment]
            out = _output_path(path, "converted", f".{fmt}")
            img.save(out, pil_fmt)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "compress":
        quality = int(params.get("quality", 70))
        try:
            img = Image.open(path).convert("RGB")  # type: ignore[assignment]
            out = _output_path(path, f"compressed_q{quality}", ".jpg")
            img.save(out, "JPEG", quality=quality, optimize=True)
            before = _file_size_str(path)
            after = _file_size_str(out)
            return f"Compressed: {before} -> {after}. Saved: {out.name}"
        except Exception as e:
            return f"Compress failed: {e}"

    if action == "crop":
        try:
            left = int(params.get("left", 0))
            top = int(params.get("top", 0))
            right = int(params.get("right", 0))
            bottom = int(params.get("bottom", 0))
            if not all([left, top, right, bottom]):
                return "Specify left, top, right, bottom crop coordinates."
            img = Image.open(path)
            cropped = img.crop((left, top, right, bottom))
            out = _output_path(path, f"cropped_{left}x{top}_{right}x{bottom}")
            cropped.save(out)
            return f"Cropped to {right - left}x{bottom - top}px. Saved: {out.name}"
        except Exception as e:
            return f"Crop failed: {e}"

    if action == "info":
        try:
            img = Image.open(path)
            return (
                f"Image: {path.name}, Format: {img.format}, "
                f"Size: {img.size[0]}x{img.size[1]}px, "
                f"Mode: {img.mode}, File size: {_file_size_str(path)}"
            )
        except Exception as e:
            return f"Info failed: {e}"

    return _process_image(path, "info", params)


# ── PDF operations ──────────────────────────────────────────────────

def _extract_pdf_text(path: Path, max_chars: int = 50000) -> str:
    text = ""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text += (page.extract_text() or "") + "\n"
    except ImportError:
        try:
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:  # type: ignore[assignment]
                    text += page.extract_text() + "\n"
        except ImportError:
            return ""
    return text[:max_chars]


def _process_pdf(path: Path, action: str, params: dict) -> str:
    action = action or "info"

    if action in ("summarize", "analyze", "reformat", "translate_hint"):
        text = _extract_pdf_text(path)
        if not text.strip():
            return "Could not extract text from PDF (may be scanned/image-based)."
        # Return extracted text — the LLM handles the analysis
        return (
            f"PDF text extracted ({len(text)} chars). "
            f"You requested: {action}. Here is the raw text for your analysis:\n\n"
            f"{text[:3000]}"
            + ("\n\n...(truncated)" if len(text) > 3000 else "")
        )

    if action in ("extract_text",):
        text = _extract_pdf_text(path)
        if not text.strip():
            return "Could not extract text from PDF (may be scanned/image-based)."
        out = _output_path(path, "text", ".txt")
        out.write_text(text, encoding="utf-8")
        return f"Text extracted ({len(text)} chars). Saved: {out.name}"

    if action == "info":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
            return f"PDF: {pages} pages, size: {_file_size_str(path)}"
        except ImportError:
            try:
                import PyPDF2
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    pages = len(reader.pages)
                return f"PDF: {pages} pages, size: {_file_size_str(path)}"
            except ImportError:
                return f"PDF size: {_file_size_str(path)}"

    if action == "to_word":
        text = _extract_pdf_text(path)
        if not text.strip():
            return "Could not extract text to convert."
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "converted", ".docx")
            doc.save(out)  # type: ignore[arg-type]
            return f"PDF converted to Word ({len(text)} chars). Saved: {out.name}"
        except ImportError:
            return "python-docx is not installed. Run: pip install python-docx"
        except Exception as e:
            return f"Word conversion failed: {e}"

    if action == "extract_pages":
        pages_str = params.get("pages", "")
        if not pages_str:
            return "Specify pages, e.g. '1,3,5' or '1-5'."
        try:
            import PyPDF2
            ranges: list[int] = []
            for part in pages_str.split(","):
                if "-" in part:
                    a, b = part.split("-")
                    ranges.extend(range(int(a) - 1, int(b)))
                else:
                    ranges.append(int(part) - 1)
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                writer = PyPDF2.PdfWriter()
                for i in ranges:
                    if i < len(reader.pages):
                        writer.add_page(reader.pages[i])
                out = _output_path(path, f"pages_{pages_str.replace(',','-')}")
                with open(out, "wb") as out_f:
                    writer.write(out_f)
                return f"Extracted {len(ranges)} page(s). Saved: {out.name}"
        except ImportError:
            return "PyPDF2 is not installed. Run: pip install PyPDF2"
        except Exception as e:
            return f"Page extraction failed: {e}"

    return f"Unknown PDF action: {action}"


# ── DOCX operations ─────────────────────────────────────────────────

def _process_docx(path: Path, action: str, _params: dict) -> str:
    action = action or "info"
    try:
        from docx import Document
        doc = Document(path)  # type: ignore[arg-type]
        text = "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        return "python-docx is not installed. Run: pip install python-docx"
    except Exception as e:
        return f"Failed to open docx: {e}"

    if action in ("summarize", "analyze", "reformat", "translate_hint"):
        if not text.strip():
            return "Document appears empty."
        return (
            f"Document text extracted ({len(text)} chars). "
            f"You requested: {action}. Raw text:\n\n{text[:3000]}"
            + ("\n\n...(truncated)" if len(text) > 3000 else "")
        )

    if action in ("extract_text",):
        if not text.strip():
            return "Document appears empty."
        out = _output_path(path, "text", ".txt")
        out.write_text(text, encoding="utf-8")
        return f"Text extracted ({len(text)} chars). Saved: {out.name}"

    if action == "info":
        return f"DOCX: {path.name}, {len(doc.paragraphs)} paragraphs, size: {_file_size_str(path)}"

    return f"Unknown DOCX action: {action}"


# ── Text operations ─────────────────────────────────────────────────

def _process_text(path: Path, action: str, _params: dict) -> str:
    action = action or "info"
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            text = path.read_text(encoding="latin-1")
        except Exception as e:
            return f"Failed to read file: {e}"

    if action in ("summarize", "analyze", "reformat", "translate_hint"):
        return (
            f"Text file content ({len(text)} chars, {len(text.splitlines())} lines). "
            f"You requested: {action}. Raw content:\n\n{text[:3000]}"
            + ("\n\n...(truncated)" if len(text) > 3000 else "")
        )

    if action == "word_count":
        words = len(re.findall(r"\b\w+\b", text))
        lines = len(text.splitlines())
        chars = len(text)
        return f"Word count: {words}, Lines: {lines}, Characters: {chars}"

    if action == "info":
        words = len(re.findall(r"\b\w+\b", text))
        lines = len(text.splitlines())
        return f"Text: {path.name}, {lines} lines, {words} words, size: {_file_size_str(path)}"

    return f"Unknown text action: {action}"


# ── CSV operations ──────────────────────────────────────────────────

def _process_csv(path: Path, action: str, params: dict) -> str:
    action = action or "info"
    import csv

    try:
        with open(path, encoding="utf-8") as f:
            reader = csv.DictReader(f)
            rows = list(reader)
            fieldnames = reader.fieldnames or []
    except Exception as e:
        return f"Failed to read CSV: {e}"

    if not rows:
        return "CSV file is empty."

    if action == "stats":
        return f"CSV: {len(rows)} rows, {len(fieldnames)} columns: {', '.join(fieldnames)}"

    if action == "info":
        return f"CSV: {path.name}, {len(rows)} rows, {len(fieldnames)} columns, size: {_file_size_str(path)}"

    if action == "filter":
        column = params.get("column", "")
        value = params.get("value", "")
        if not column or not value:
            return "Specify column and value to filter by."
        filtered = [r for r in rows if r.get(column, "").lower() == value.lower()]
        return f"Filtered: {len(filtered)} of {len(rows)} rows match {column}={value}."

    if action == "sort":
        column = params.get("column", "")
        reverse = params.get("reverse", False)
        if not column:
            return "Specify column to sort by."
        sorted_rows = sorted(rows, key=lambda r: r.get(column, ""), reverse=reverse)
        preview = "\n".join(
            f"{r.get(fieldnames[0], '')}, {r.get(fieldnames[1] if len(fieldnames) > 1 else fieldnames[0], '')}"
            for r in sorted_rows[:5]
        )
        return f"Sorted by {column} ({'desc' if reverse else 'asc'}). Preview:\n{preview}"

    if action == "convert":
        fmt = params.get("format", "json").lower()
        if fmt == "json":
            out = _output_path(path, "converted", ".json")
            out.write_text(json.dumps(rows, indent=2), encoding="utf-8")
            return f"CSV converted to JSON ({len(rows)} records). Saved: {out.name}"
        return f"Unsupported format: {fmt}. Supported: json"

    return f"Unknown CSV action: {action}"


# ── Excel operations ────────────────────────────────────────────────

def _process_excel(path: Path, action: str, params: dict) -> str:
    action = action or "info"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except ImportError:
        return "openpyxl is not installed. Run: pip install openpyxl"
    except Exception as e:
        return f"Failed to open Excel file: {e}"

    sheets = wb.sheetnames

    if action == "info":
        data = {s: wb[s].max_row for s in sheets}
        return f"Excel: {path.name}, {len(sheets)} sheet(s): {', '.join(f'{s} ({data[s]} rows)' for s in sheets)}"

    if action in ("analyze",):
        # Return structure for LLM analysis
        sheet_name = params.get("sheet", sheets[0])
        ws = wb[sheet_name]
        headers = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1), [])]
        return (
            f"Excel sheet '{sheet_name}': {ws.max_row} rows, "
            f"columns: {', '.join(str(h) for h in headers if h)}. "
            f"The LLM can analyze this data."
        )

    if action == "convert":
        fmt = params.get("format", "csv").lower()
        sheet_name = params.get("sheet", sheets[0])
        ws = wb[sheet_name]
        if fmt == "csv":
            import csv
            out = _output_path(path, f"{sheet_name}_converted", ".csv")
            with open(out, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow(row)
            return f"Sheet '{sheet_name}' converted to CSV. Saved: {out.name}"
        return f"Unsupported format: {fmt}"

    return f"Unknown Excel action: {action}"


# ── JSON operations ─────────────────────────────────────────────────

def _process_json(path: Path, action: str, params: dict) -> str:
    action = action or "info"
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception as e:
        return f"Failed to parse JSON: {e}"

    if action == "validate":
        return "Valid JSON." if isinstance(data, (dict, list)) else f"Valid JSON ({type(data).__name__})."

    if action == "format":
        indent = int(params.get("indent", 2))
        formatted = json.dumps(data, indent=indent)
        out = _output_path(path, "formatted", ".json")
        out.write_text(formatted, encoding="utf-8")
        return f"Formatted JSON. Saved: {out.name}"

    if action == "extract":
        key = params.get("key", "")
        if not key:
            return "Specify key to extract."
        if isinstance(data, dict):
            value = json.dumps(data.get(key, "Key not found"), indent=2)
        else:
            return "JSON root is not an object."
        return f"Extracted key '{key}':\n{value[:2000]}"

    if action == "info":
        if isinstance(data, dict):
            return f"JSON object: {len(data)} keys, size: {_file_size_str(path)}"
        if isinstance(data, list):
            return f"JSON array: {len(data)} items, size: {_file_size_str(path)}"
        return f"JSON: {type(data).__name__}, size: {_file_size_str(path)}"

    return f"Unknown JSON action: {action}"


# ── Code operations ─────────────────────────────────────────────────

def _process_code(path: Path, action: str, _params: dict) -> str:
    action = action or "info"
    try:
        text = path.read_text(encoding="utf-8")
    except Exception as e:
        return f"Failed to read file: {e}"

    if action in ("explain", "review", "fix", "document"):
        return (
            f"Code file: {path.name} ({len(text.splitlines())} lines, {len(text)} chars). "
            f"You requested: {action}. Raw code:\n\n{text[:3000]}"
            + ("\n\n...(truncated)" if len(text) > 3000 else "")
        )

    if action == "word_count":
        words = len(re.findall(r"\b\w+\b", text))
        lines = len(text.splitlines())
        chars = len(text)
        return f"Word count: {words}, Lines: {lines}, Characters: {chars}"

    if action == "run":
        try:
            ext = path.suffix.lower()
            if ext in (".py",):
                result = subprocess.run(
                    ["python", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                output = result.stdout or result.stderr
                return f"Output:\n{output[:2000]}" if output else "Script completed with no output."
            return f"Running {ext} files is not supported. Use the run_command tool instead."
        except subprocess.TimeoutExpired:
            return "Script timed out after 30 seconds."
        except Exception as e:
            return f"Failed to run: {e}"

    if action == "info":
        lines = text.splitlines()  # type: ignore[assignment]
        return f"Code: {path.name}, {len(lines)} lines, {len(text)} chars, size: {_file_size_str(path)}"  # type: ignore[arg-type]

    return f"Unknown code action: {action}"


# ── Audio operations ────────────────────────────────────────────────

def _process_audio(path: Path, action: str, _params: dict) -> str:
    action = action or "info"

    if action == "info":
        try:
            import mutagen
            audio = mutagen.File(path)
            if audio:
                info = audio.info
                length = getattr(info, "length", 0)
                return f"Audio: {path.name}, {length:.1f}s, size: {_file_size_str(path)}"
        except ImportError:
            pass
        return f"Audio: {path.name}, size: {_file_size_str(path)}"

    if action in ("transcribe",):
        return "Speech-to-text transcription is handled by the STT module. Use voice input instead."

    return f"Unknown audio action: {action}"


# ── Video operations ────────────────────────────────────────────────

def _process_video(path: Path, action: str, _params: dict) -> str:
    action = action or "info"

    if action == "info":
        try:
            import importlib
            ffmpeg = importlib.import_module("ffmpeg")
            probe = ffmpeg.probe(str(path))
            fmt = probe.get("format", {})
            duration = float(fmt.get("duration", 0))
            return f"Video: {path.name}, {duration:.1f}s, size: {_file_size_str(path)}"
        except ImportError:
            return f"Video: {path.name}, size: {_file_size_str(path)}"
        except Exception:
            return f"Video: {path.name}, size: {_file_size_str(path)}"

    if action == "extract_audio":
        try:
            import importlib
            ffmpeg = importlib.import_module("ffmpeg")
            out = _output_path(path, "audio", ".mp3")
            ffmpeg.input(str(path)).output(str(out), acodec="libmp3lame", vn=None).run(quiet=True, overwrite_output=True)
            return f"Audio extracted. Saved: {out.name}"
        except ImportError:
            return "ffmpeg-python is not installed. Run: pip install ffmpeg-python"
        except Exception as e:
            return f"Audio extraction failed: {e}"

    return f"Unknown video action: {action}"


# ── Archive operations ──────────────────────────────────────────────

def _process_archive(path: Path, action: str, params: dict) -> str:
    action = action or "info"

    if action in ("list", "info"):
        import zipfile
        try:
            with zipfile.ZipFile(path, "r") as zf:
                names = zf.namelist()
                return f"Archive: {path.name}, {len(names)} file(s), size: {_file_size_str(path)}" + \
                    ("\nContents:\n" + "\n".join(names[:20]) + ("\n..." if len(names) > 20 else "") if action == "list" else "")
        except zipfile.BadZipFile:
            return f"Archive: {path.name}, size: {_file_size_str(path)} (not a zip, try extracting manually)"
        except Exception as e:
            return f"Failed to read archive: {e}"

    if action == "extract":
        import zipfile
        dest = params.get("dest", str(path.parent / path.stem))
        try:
            with zipfile.ZipFile(path, "r") as zf:
                zf.extractall(dest)
            return f"Extracted to: {dest}"
        except zipfile.BadZipFile:
            return "Not a valid zip file. Other formats need 7-Zip or similar."
        except Exception as e:
            return f"Extraction failed: {e}"

    return f"Unknown archive action: {action}"


# ── PPTX operations ─────────────────────────────────────────────────

def _process_pptx(path: Path, action: str, _params: dict) -> str:
    action = action or "info"
    try:
        from pptx import Presentation
        prs = Presentation(path)  # type: ignore[arg-type]
    except ImportError:
        return "python-pptx is not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"Failed to open PPTX: {e}"

    if action in ("summarize", "analyze"):
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"Slide {i}: {' | '.join(slide_text)}")
        content = "\n".join(text_parts)
        return (
            f"Presentation: {len(prs.slides)} slides. "
            f"You requested: {action}. Content:\n\n{content[:3000]}"
            + ("\n\n...(truncated)" if len(content) > 3000 else "")
        )

    if action == "extract_text":
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        text = "\n\n".join(text_parts)
        out = _output_path(path, "text", ".txt")
        out.write_text(text, encoding="utf-8")
        return f"Text extracted from {len(prs.slides)} slides ({len(text)} chars). Saved: {out.name}"

    if action == "info":
        return f"PPTX: {path.name}, {len(prs.slides)} slides, size: {_file_size_str(path)}"

    return f"Unknown PPTX action: {action}"


# ── Generic handler ─────────────────────────────────────────────────

def _process_unknown(path: Path, action: str, params: dict) -> str:
    return f"Unknown file type: {path.suffix}. Supported types: image, pdf, docx, text, csv, excel, json, code, audio, video, archive, pptx."


# ── Dispatch ────────────────────────────────────────────────────────

_ACTION_MAP = {
    "image":   _process_image,
    "pdf":     _process_pdf,
    "docx":    _process_docx,
    "text":    _process_text,
    "csv":     _process_csv,
    "excel":   _process_excel,
    "json":    _process_json,
    "code":    _process_code,
    "audio":   _process_audio,
    "video":   _process_video,
    "archive": _process_archive,
    "pptx":    _process_pptx,
}


# ── Public API ──────────────────────────────────────────────────────

def process_file(file_path: str, action: str = "info", **params) -> str:
    """
    Process a file with the given action.

    Args:
        file_path: Absolute or relative path to the file.
        action: Operation to perform (e.g. 'info', 'resize', 'convert', 'extract_text').
        **params: Additional parameters specific to the action.

    Returns:
        Result string describing what was done.
    """
    path = Path(file_path).resolve()
    if not path.exists():
        return f"File not found: {file_path}"

    if not path.is_file():
        return f"Not a file: {file_path}"

    file_type = _detect_type(path)
    handler = _ACTION_MAP.get(file_type, _process_unknown)
    logger.info("FileProcessor: %s: %s action=%s", file_type, path.name, action)
    return handler(path, action, params)

